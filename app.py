import streamlit as st
# Build: PU ExpertCenter Enterprise - Internal Version
import openai
import os
from pathlib import Path
import PyPDF2
from docx import Document
from dotenv import load_dotenv
import re
from collections import Counter
import json
import time
from datetime import datetime
from PIL import Image

# Load environment variables
load_dotenv()

# Flexible resolver for documents directory (accepts multiple name variants)
from typing import Optional

def resolve_documents_dir() -> Optional[str]:
    base_dir = Path(__file__).parent
    # Look in current dir first
    candidates = [p for p in base_dir.iterdir() if p.is_dir()]
    # Also consider repo root (when app is in a subfolder on Streamlit)
    repo_root = Path('.')
    try:
        candidates += [p for p in repo_root.iterdir() if p.is_dir()]
    except Exception:
        pass
    acceptable = {"document database", "documents database"}
    for p in candidates:
        try:
            if p.name.lower() in acceptable:
                return str(p)
        except Exception:
            continue
    # Fallback to conventional path relative to file
    fallback = base_dir / "Document Database"
    return str(fallback) if fallback.exists() else None

class PUExpertCenterMinimal:
    def __init__(self):
        # Try Streamlit secrets first, fallback to environment variables
        try:
            api_key = st.secrets["OPENAI_API_KEY"]
            assistant_id = st.secrets.get("OPENAI_ASSISTANT_ID", "")
        except:
            # Fallback to environment variables for local development
            api_key = os.getenv("OPENAI_API_KEY")
            assistant_id = os.getenv("OPENAI_ASSISTANT_ID", "")
        
        self.openai_client = openai.OpenAI(api_key=api_key)
        self.documents = []
        self.processed = False
        # Assistant (OpenAI Assistants API)
        self.assistant_id = assistant_id
        # Resume/cache/tracking
        self.cache_path = Path(__file__).parent / "processed_cache.json"
        self.chunks_store_path = Path(__file__).parent / "chunks_store.jsonl"
        self.progress_path = Path(__file__).parent / ".progress.json"
        self.log_file = Path(__file__).parent / "processing_log.txt"
        self.log_entries = []
        self.lock_file = Path(__file__).parent / "app_lock.json"
        self._processed_index = {}  # file_path -> {mtime, size}
        self._load_cache_and_chunks()
        
    def extract_text_from_file(self, file_path):
        """Extract text from various file formats"""
        try:
            file_path = Path(file_path)
            if file_path.suffix.lower() == '.pdf':
                # Try multiple PDF extraction methods
                text = ""
                
                # Method 1: Try PyMuPDF first (often better than PyPDF2)
                try:
                    import fitz
                    doc = fitz.open(file_path)
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        text += page.get_text() + "\n"
                    doc.close()
                    if text.strip():
                        return text
                except Exception as e:
                    print(f"PyMuPDF failed for {file_path.name}: {e}")
                
                # Method 2: Try PyPDF2
                try:
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        text = ""
                        for page in pdf_reader.pages:
                            text += page.extract_text() + "\n"
                    if text.strip():
                        return text
                except Exception as e:
                    print(f"PyPDF2 failed for {file_path.name}: {e}")
                
                # Method 3: Try OCR as last resort
                try:
                    import pytesseract
                    from PIL import Image
                    import fitz
                    
                    doc = fitz.open(file_path)
                    ocr_text = ""
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        pix = page.get_pixmap()
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        ocr_text += pytesseract.image_to_string(img) + "\n"
                    doc.close()
                    if ocr_text.strip():
                        return ocr_text
                except Exception as e:
                    print(f"OCR failed for {file_path.name}: {e}")
                
                # If all methods failed, return a message indicating the issue
                return f"[PDF file: {file_path.name} - All extraction methods failed. File may be image-based or corrupted.]"
            elif file_path.suffix.lower() == '.docx':
                # Extract paragraphs and table text; then fallback to docx2txt if empty
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    if paragraph.text:
                        text += paragraph.text + "\n"
                # Include table cell text (often missed)
                for table in getattr(doc, 'tables', []):
                    for row in table.rows:
                        row_text = " ".join(cell.text for cell in row.cells if cell.text)
                        if row_text.strip():
                            text += row_text + "\n"
                if text.strip():
                    return text
                # Fallback: docx2txt can sometimes recover more content
                try:
                    import docx2txt
                    fallback = docx2txt.process(str(file_path))
                    return fallback if fallback and fallback.strip() else ""
                except Exception:
                    return ""
            elif file_path.suffix.lower() == '.txt':
                with open(file_path, 'r', encoding='utf-8') as file:
                    return file.read()
            elif file_path.suffix.lower() == '.csv':
                # Read CSV as plain text (simple, dependency-free)
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            elif file_path.suffix.lower() in ['.md']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            elif file_path.suffix.lower() in ['.html', '.htm']:
                # Naive HTML to text: strip tags
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    raw = f.read()
                text = re.sub(r'<script[\s\S]*?</script>', ' ', raw, flags=re.IGNORECASE)
                text = re.sub(r'<style[\s\S]*?</style>', ' ', text, flags=re.IGNORECASE)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text)
                return text
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                # Extract text from PowerPoint files
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                except ImportError:
                    return f"[PowerPoint file: {file_path.name} - python-pptx not installed]"
            elif file_path.suffix.lower() in ['.xlsx', '.xls']:
                # Extract text from Excel files
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path)
                    text = ""
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " ".join(str(cell) for cell in row if cell is not None)
                            if row_text.strip():
                                text += row_text + "\n"
                    return text
                except ImportError:
                    return f"[Excel file: {file_path.name} - openpyxl not installed]"
            elif file_path.suffix.lower() == '.doc':
                # Extract text from legacy Word files
                try:
                    import docx2txt
                    return docx2txt.process(file_path)
                except ImportError:
                    return f"[Legacy Word file: {file_path.name} - python-docx2txt not installed]"
            else:
                return ""
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""
    
    def chunk_text(self, text, chunk_size=1000, overlap=200):
        """Split text into overlapping chunks for better retrieval"""
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start = end - overlap
        return chunks
    
    def process_documents(self, documents_folder):
        """Process all documents in the folder with resume, ETA, and progress tracking"""
        st.write("🔄 Processing PU documents...")
        start_time = time.time()
        
        # Initialize logging
        self._log("info", f"Starting document processing for: {documents_folder}")
        failed_files = {}
        skipped_files = {}
        successful_files = []
        processed_count = 0
        
        # Discover candidate files recursively (exclude program dir)
        base_path = Path(documents_folder)
        program_dir_name = 'PU_ExpertCenter_RAG'
        patterns = ['**/*.pdf', '**/*.docx', '**/*.txt', '**/*.csv', '**/*.md', '**/*.html', '**/*.htm', '**/*.pptx', '**/*.xlsx', '**/*.xls', '**/*.ppt', '**/*.doc']
        all_files = []
        for pattern in patterns:
            all_files.extend(base_path.rglob(pattern))
        # Skip temp/lock files like ~$*.docx and anything under program dir
        doc_files = []
        for p in all_files:
            if p.name.startswith('~$'):
                continue
            parts = set(p.parts)
            if program_dir_name in parts:
                continue
            doc_files.append(p)
        total_files = len(doc_files)
        processed = 0
        
        # Load prior documents so resume works
        existing_docs = len(self.documents)
        progress_bar = st.progress(0.0)
        current_file_placeholder = st.empty()
        eta_placeholder = st.empty()
        count_placeholder = st.empty()
        
        for file_path in doc_files:
            try:
                stat = file_path.stat()
                key = str(file_path.resolve())
                mtime = int(stat.st_mtime)
                size = stat.st_size
                already = self._processed_index.get(key)
                # If unchanged and already processed, skip
                if already and already.get('mtime') == mtime and already.get('size') == size:
                    processed += 1
                    self._write_progress(processed, total_files, file_path.name, start_time)
                    self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
                    continue
            except Exception:
                pass
            
            current_file_placeholder.write(f"Processing: {file_path.name}")
            
            # Extract text
            try:
                text = self.extract_text_from_file(file_path)
                if not text.strip():
                    self._log("warning", "No text extracted from file", str(file_path))
                    skipped_files[str(file_path)] = "No text content"
                    processed += 1
                    self._write_progress(processed, total_files, file_path.name, start_time)
                    self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
                    continue
            except Exception as e:
                self._log("error", f"Failed to extract text: {str(e)}", str(file_path))
                failed_files[str(file_path)] = str(e)
                processed += 1
                self._write_progress(processed, total_files, file_path.name, start_time)
                self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
                continue
            
            # Remove any previously stored segments for this file in memory
            self._remove_existing_file_chunks_in_memory(str(file_path))
            
            # Create document segments
            chunks = self.chunk_text(text)
            
            # Add to documents list and persist to chunks_store
            self._append_chunks(file_path, chunks)
            
            # Log successful processing
            self._log("info", f"Successfully processed {len(chunks)} segments", str(file_path))
            successful_files.append(str(file_path))
            processed_count += 1
            
            # Update cache index for resume
            self._processed_index[str(file_path.resolve())] = {'mtime': int(file_path.stat().st_mtime), 'size': file_path.stat().st_size}
            self._save_cache()
            
            processed += 1
            self._write_progress(processed, total_files, file_path.name, start_time)
            self._update_progress_ui(progress_bar, current_file_placeholder, eta_placeholder, count_placeholder, processed, total_files, file_path.name, start_time)
        
        self.processed = True
        self._finalize_progress()
        
        # Write processing summary to log
        self._write_log_summary(total_files, processed_count, failed_files, skipped_files, successful_files)
        self._log("info", f"Processing completed. Successfully processed: {processed_count}, Failed: {len(failed_files)}, Skipped: {len(skipped_files)}")
        
        st.success(f"✅ Processed {processed} of {total_files} files; total documents loaded: {len(self.documents)} (prev: {existing_docs})")
    
    def search_documents(self, query, n_results=5):
        """Search for relevant segments using token overlap + fuzzy matching (robust)."""
        if not self.processed:
            return []

        try:
            from rapidfuzz import fuzz
        except Exception:
            fuzz = None

        # Basic stopwords to improve signal
        stopwords = {
            'the','a','an','and','or','of','to','in','for','on','at','by','with','as','is','are','be','was','were','it','this','that','these','those','from','about','into','over','under','between','within','per'
        }
        q_lower = query.lower()
        query_tokens = [t for t in re.split(r"[^a-z0-9]+", q_lower) if t and t not in stopwords]
        if not query_tokens:
            query_tokens = [t for t in re.split(r"[^a-z0-9]+", q_lower) if t]

        results = []
        for doc in self.documents:
            text = doc['text']
            words = doc['words']
            filename = doc.get('filename','')

            # Token overlap score
            word_freq = Counter(words)
            overlap = sum(word_freq.get(tok, 0) for tok in query_tokens) / max(1, len(words))

            # Fuzzy score (partial ratio) on text
            fuzzy_score = 0.0
            if fuzz is not None:
                try:
                    fuzzy_score = fuzz.partial_ratio(q_lower, text.lower()) / 100.0
                except Exception:
                    fuzzy_score = 0.0

            # Filename boost if query mentions terms present in filename
            filename_boost = 0.0
            fname_lower = filename.lower()
            if any(tok in fname_lower for tok in query_tokens):
                filename_boost = 0.1

            # Combined score
            score = 0.6 * overlap + 0.4 * fuzzy_score + filename_boost

            if score > 0.01:
                matched_words = [tok for tok in set(query_tokens) if tok in words or tok in fname_lower]
                results.append({
                    'text': text,
                    'filename': filename,
                    'file_path': doc['file_path'],
                    'similarity': score,
                    'matched_words': matched_words
                })

        results.sort(key=lambda x: x['similarity'], reverse=True)
        return results[:n_results]
    
    def generate_answer(self, query, context_chunks):
        """Generate answer using OpenAI with retrieved context"""
        context = "\n\n".join([chunk['text'] for chunk in context_chunks])
        
        prompt = f"""
        You are a polyurethane expert with access to a comprehensive knowledge base of PU research documents.
        
        Question: {query}
        
        Context from research documents:
        {context}
        
        Please provide a comprehensive, well-structured answer based on the provided context. 
        Include specific details, data points, and cite the source documents when possible.
        If the context doesn't contain enough information to fully answer the question, 
        clearly state what information is missing.
        
        Format your response as a professional PU expert report with:
        1. Executive Summary
        2. Key Findings
        3. Supporting Evidence
        4. Source References
        """
        
        try:
            response = self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a polyurethane expert providing detailed, accurate information based on research documents."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=2000,
                temperature=0.3
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error generating answer: {e}"
    
    def answer_question(self, question):
        """Main method to answer a question using Assistant if configured, else local retrieval"""
        # Prefer Assistant (Vector Store + WebSearch) when configured
        if self.assistant_id:
            try:
                thread = self.openai_client.beta.threads.create()
                self.openai_client.beta.threads.messages.create(
                    thread_id=thread.id,
                    role="user",
                    content=question
                )
                run = self.openai_client.beta.threads.runs.create(
                    thread_id=thread.id,
                    assistant_id=self.assistant_id
                )
                # Poll until completion
                start = time.time()
                while True:
                    r = self.openai_client.beta.threads.runs.retrieve(thread_id=thread.id, run_id=run.id)
                    if r.status in ["completed", "failed", "cancelled", "expired"]:
                        break
                    time.sleep(0.4)
                    # optional safety timeout
                    if time.time() - start > 120:
                        return "Assistant timeout. Please try again.", []
                if r.status != "completed":
                    return f"Assistant run status: {r.status}", []
                msgs = self.openai_client.beta.threads.messages.list(thread_id=thread.id)
                answer_text = ""
                sources = []
                for m in msgs.data:
                    if m.role == 'assistant':
                        parts = []
                        for c in m.content:
                            if getattr(c, 'type', '') == 'text':
                                parts.append(c.text.value)
                        if parts:
                            answer_text = "\n\n".join(parts)
                        # If the assistant returns file references as attachments
                        if hasattr(m, 'attachments') and m.attachments:
                            for att in m.attachments:
                                fname = getattr(att, 'filename', None)
                                if fname:
                                    sources.append({'filename': fname, 'similarity_score': 1.0})
                        break
                
                # Extract source references from the answer text if no attachments found
                if not sources and answer_text:
                    import re
                    # Look for file references in the text (common patterns)
                    file_refs = re.findall(r'\[([^\]]+\.(?:pdf|docx?|txt|pptx?|xlsx?|html?|md|csv))\]', answer_text, re.IGNORECASE)
                    for ref in file_refs:
                        sources.append({'filename': ref, 'similarity_score': 1.0})
                    
                    # Also look for quoted filenames
                    quoted_files = re.findall(r'"([^"]+\.(?:pdf|docx?|txt|pptx?|xlsx?|html?|md|csv))"', answer_text, re.IGNORECASE)
                    for ref in quoted_files:
                        sources.append({'filename': ref, 'similarity_score': 1.0})
                
                return (answer_text or ""), sources
            except Exception as e:
                return f"Assistant error: {e}", []
        
        # Fallback to local retrieval
        search_results = self.search_documents(question)
        if not search_results:
            return "No relevant information found in the knowledge base.", []
        answer = self.generate_answer(question, search_results)
        sources = []
        for result in search_results:
            sources.append({'filename': result['filename'], 'similarity_score': result['similarity'], 'matched_words': result['matched_words']})
        return answer, sources

    # -------------- Internal helpers for cache/progress --------------
    def _load_cache_and_chunks(self):
        try:
            if self.cache_path.exists():
                with open(self.cache_path, 'r', encoding='utf-8') as f:
                    self._processed_index = json.load(f)
            else:
                self._processed_index = {}
        except Exception:
            self._processed_index = {}
        
        # Load previously stored document segments into memory
        try:
            if self.chunks_store_path.exists():
                with open(self.chunks_store_path, 'r', encoding='utf-8') as f:
                    for line in f:
                        try:
                            rec = json.loads(line)
                            self.documents.append(rec)
                        except Exception:
                            continue
                if self.documents:
                    self.processed = True
        except Exception:
            pass
    
    def _read_progress(self) -> dict:
        try:
            if self.progress_path.exists():
                with open(self.progress_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception:
            pass
        return {}

    def _get_recent_log_lines(self, max_lines: int = 50) -> list:
        try:
            if self.log_file.exists():
                with open(self.log_file, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                return lines[-max_lines:]
        except Exception:
            pass
        return []

    def _clear_cache_and_chunks(self):
        try:
            if self.cache_path.exists():
                self.cache_path.unlink()
        except Exception:
            pass
        try:
            if self.chunks_store_path.exists():
                self.chunks_store_path.unlink()
        except Exception:
            pass
        self._processed_index = {}
        self.documents = []
        self.processed = False
    
    def _save_cache(self):
        try:
            with open(self.cache_path, 'w', encoding='utf-8') as f:
                json.dump(self._processed_index, f)
        except Exception:
            pass
    
    def _append_chunks(self, file_path: Path, chunks):
        records = []
        for i, chunk in enumerate(chunks):
            if not chunk.strip():
                continue
            rec = {
                'filename': file_path.name,
                'file_path': str(file_path),
                'chunk_id': i,
                'text': chunk,
                'file_type': file_path.suffix,
                'words': chunk.lower().split()
            }
            self.documents.append(rec)
            records.append(rec)
        # Append to persistent store
        try:
            with open(self.chunks_store_path, 'a', encoding='utf-8') as f:
                for rec in records:
                    f.write(json.dumps(rec, ensure_ascii=False) + "\n")
        except Exception:
            pass
    
    def _remove_existing_file_chunks_in_memory(self, file_path_str: str):
        if not self.documents:
            return
        self.documents = [d for d in self.documents if d.get('file_path') != file_path_str]
        # Rebuild chunks_store without this file (simple rewrite)
        try:
            if self.chunks_store_path.exists():
                with open(self.chunks_store_path, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                with open(self.chunks_store_path, 'w', encoding='utf-8') as f:
                    for line in lines:
                        try:
                            rec = json.loads(line)
                            if rec.get('file_path') != file_path_str:
                                f.write(line)
                        except Exception:
                            continue
        except Exception:
            pass
    
    def _write_progress(self, processed, total, current_file, start_time):
        try:
            elapsed = max(0.001, time.time() - start_time)
            pct = (processed / total) if total else 1.0
            eta_sec = (elapsed / processed) * (total - processed) if processed > 0 and total > 0 else 0
            data = {
                'processed': processed,
                'total': total,
                'percent': round(pct * 100, 2),
                'current_file': current_file,
                'elapsed_seconds': int(elapsed),
                'eta_seconds': int(eta_sec)
            }
            with open(self.progress_path, 'w', encoding='utf-8') as f:
                json.dump(data, f)
        except Exception:
            pass
    
    def _update_progress_ui(self, bar, file_ph, eta_ph, count_ph, processed, total, current_file, start_time):
        pct = (processed / total) if total else 1.0
        bar.progress(pct)
        elapsed = max(0.001, time.time() - start_time)
        eta_sec = int((elapsed / processed) * (total - processed)) if processed > 0 and total > 0 else 0
        mins, secs = divmod(max(0, eta_sec), 60)
        eta_ph.write(f"ETA: {mins}m {secs}s remaining")
        count_ph.write(f"Files: {processed}/{total} | Current: {current_file}")
    
    def _finalize_progress(self):
        try:
            if self.progress_path.exists():
                with open(self.progress_path, 'w', encoding='utf-8') as f:
                    json.dump({'status': 'complete'}, f)
        except Exception:
            pass
    
    def _get_document_files(self, doc_dir: str) -> list:
        """Get list of all document files to process"""
        if not os.path.exists(doc_dir):
            return []
        
        doc_files = []
        for ext in ['.pdf', '.docx', '.txt', '.csv', '.md', '.html', '.htm', '.pptx', '.xlsx', '.xls', '.ppt', '.doc']:
            for file_path in Path(doc_dir).rglob(f'*{ext}'):
                if 'PU_ExpertCenter_RAG' not in str(file_path) and not file_path.name.startswith('~$'):
                    doc_files.append(str(file_path))
        return doc_files
    
    def check_for_new_files(self, doc_dir: str) -> int:
        """Check for new or updated files without processing them (excludes failed and skipped files)"""
        if not os.path.exists(doc_dir):
            return 0
        
        all_files = self._get_document_files(doc_dir)
        failed_files = self._get_failed_files()
        skipped_files = self._get_skipped_files()
        
        # Filter out failed and skipped files from the count
        valid_files = []
        for file_path in all_files:
            file_name = Path(file_path).name
            if file_name not in failed_files and file_name not in skipped_files:
                valid_files.append(file_path)
        
        return len(valid_files)
    
    def _log(self, level, message, file_path=None):
        """Add entry to processing log"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        entry = {
            "timestamp": timestamp,
            "level": level,
            "message": message,
            "file_path": file_path
        }
        self.log_entries.append(entry)
        
        # Write to log file
        with open(self.log_file, "a", encoding="utf-8") as f:
            f.write(f"[{timestamp}] {level.upper()}: {message}")
            if file_path:
                f.write(f" | File: {file_path}")
            f.write("\n")
    
    def _write_log_summary(self, total_files, processed_files, failed_files, skipped_files, successful_files=None):
        """Write summary to log file"""
        with open(self.log_file, "a", encoding="utf-8") as f:
            f.write(f"\n{'='*60}\n")
            f.write(f"PROCESSING SUMMARY - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"{'='*60}\n")
            f.write(f"Total files found: {total_files}\n")
            f.write(f"Successfully processed: {processed_files}\n")
            f.write(f"Failed to process: {len(failed_files)}\n")
            f.write(f"Skipped (duplicates/unchanged): {len(skipped_files)}\n")
            
            if successful_files:
                f.write(f"\nSUCCESSFULLY PROCESSED FILES:\n")
                for file_path in successful_files:
                    f.write(f"  - {file_path}\n")
            
            f.write(f"\nFAILED FILES:\n")
            for file_path, reason in failed_files.items():
                f.write(f"  - {file_path}: {reason}\n")
            f.write(f"\nSKIPPED FILES:\n")
            for file_path, reason in skipped_files.items():
                f.write(f"  - {file_path}: {reason}\n")
            f.write(f"{'='*60}\n\n")
        
        # Write failed files to separate log
        self._write_failed_files_log(failed_files)
    
    def _write_failed_files_log(self, failed_files):
        """Write failed files to a separate, easily readable log"""
        failed_log_path = Path(__file__).parent / "failed_files_log.txt"
        with open(failed_log_path, "w", encoding="utf-8") as f:
            f.write(f"FAILED FILES LOG - {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"{'='*50}\n\n")
            if failed_files:
                for file_path, reason in failed_files.items():
                    f.write(f"❌ FAILED: {Path(file_path).name}\n")
                    f.write(f"   Path: {file_path}\n")
                    f.write(f"   Reason: {reason}\n")
                    f.write(f"   {'-'*40}\n")
            else:
                f.write("✅ No files failed to process.\n")
    
    def _get_failed_files(self):
        """Get list of failed files from the failed files log"""
        failed_log_path = Path(__file__).parent / "failed_files_log.txt"
        failed_files = []
        if failed_log_path.exists():
            with open(failed_log_path, "r", encoding="utf-8") as f:
                content = f.read()
                # Extract failed file paths
                import re
                failed_matches = re.findall(r"❌ FAILED: (.+?)\n", content)
                failed_files = [match.strip() for match in failed_matches]
        return failed_files
    
    def _get_skipped_files(self):
        """Get list of skipped files from the processing log"""
        skipped_files = []
        if self.log_file.exists():
            with open(self.log_file, "r", encoding="utf-8") as f:
                content = f.read()
                # Extract skipped file paths from log - look for the actual pattern in the log
                import re
                # Pattern: "  - documents database/filename.pdf: No text content"
                pattern = r"  - documents database/(.+?): No text content"
                matches = re.findall(pattern, content)
                for match in matches:
                    # Normalize to just the base filename so comparisons match
                    filename = Path(match.strip()).name
                    if filename and filename not in skipped_files:
                        skipped_files.append(filename)
        return skipped_files
    
    def _acquire_lock(self, user_name):
        """Acquire application lock to prevent conflicts"""
        try:
            if self.lock_file.exists():
                with open(self.lock_file, 'r', encoding='utf-8') as f:
                    lock_data = json.load(f)
                if time.time() - lock_data.get('timestamp', 0) > 300:
                    self._release_lock()
                else:
                    return False, lock_data.get('user', 'Unknown')
            
            lock_data = {
                'user': user_name,
                'timestamp': time.time(),
                'pid': os.getpid()
            }
            with open(self.lock_file, 'w', encoding='utf-8') as f:
                json.dump(lock_data, f)
            return True, user_name
        except Exception:
            return False, "Error"
    
    def _release_lock(self):
        """Release application lock"""
        try:
            if self.lock_file.exists():
                self.lock_file.unlink()
        except Exception:
            pass
    
    def _check_who_online(self):
        """Check who is currently using the application"""
        try:
            if self.lock_file.exists():
                with open(self.lock_file, 'r', encoding='utf-8') as f:
                    lock_data = json.load(f)
                if time.time() - lock_data.get('timestamp', 0) > 300:
                    self._release_lock()
                    return None
                return lock_data.get('user', 'Unknown')
            return None
        except Exception:
            return None

def main():
    st.set_page_config(
        page_title="PU ExpertCenter RAG",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Custom CSS to make sidebar wider
    st.markdown("""
    <style>
    .css-1d391kg {
        width: 360px !important;
        max-width: 360px !important;
    }
    .css-1lcbmhc .css-1d391kg {
        width: 360px !important;
        max-width: 360px !important;
    }
    .sidebar .sidebar-content {
        width: 360px !important;
        max-width: 360px !important;
    }
    [data-testid="stSidebar"] {
        width: 360px !important;
        min-width: 360px !important;
    }
    [data-testid="stSidebar"] > div {
        width: 360px !important;
        min-width: 360px !important;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Header with logo
    col1, col2 = st.columns([1, 4])
    with col1:
        try:
            st.image("PU ExpertCenter Logo V1.png", width=120)
        except:
            st.markdown("🧪")  # Fallback if logo not found
    with col2:
        st.title("PU ExpertCenter")
    
    st.markdown("Ask questions about polyurethane foam technology, chemistry, and applications")
    
    # Initialize RAG system
    if 'rag_system' not in st.session_state:
        st.session_state.rag_system = PUExpertCenterMinimal()
    
    # User identification and lock management
    if 'user_name' not in st.session_state:
        st.session_state.user_name = "Stefan Hermes"
    
    # Check who's processing documents
    current_user = st.session_state.user_name
    processing_user = st.session_state.rag_system._check_who_online()
    
    # Show processing status
    if processing_user and processing_user != current_user:
        st.warning(f"⚠️ {processing_user} is currently processing documents")
    elif processing_user == current_user:
        st.success(f"✅ You ({current_user}) are currently processing documents")
    
    # Sidebar for document processing and statistics
    with st.sidebar:
        st.header("📚 Document Management")
        documents_dir = resolve_documents_dir()
        if not documents_dir:
            st.error("Documents folder not found. Create 'Document Database' (or 'documents database') at repo root.")
            documents_dir = "./Document Database"  # keep UI functional
        st.caption(f"Folder: {documents_dir}")

        # Precompute file sets and exclusions
        new_files_count = st.session_state.rag_system.check_for_new_files(documents_dir)
        all_files_list = st.session_state.rag_system._get_document_files(documents_dir)
        failed_names = set(st.session_state.rag_system._get_failed_files())
        skipped_names = set(st.session_state.rag_system._get_skipped_files())
        def _is_valid(path_str: str) -> bool:
            name = Path(path_str).name
            return name not in failed_names and name not in skipped_names
        valid_files_list = [p for p in all_files_list if _is_valid(p)]
        
        if st.session_state.rag_system.processed:
            total_available = len(valid_files_list)
            processed_files = len({d.get('file_path') for d in st.session_state.rag_system.documents if d.get('file_path')})
            actual_new_files = max(0, total_available - processed_files)
        else:
            actual_new_files = new_files_count
        
        if actual_new_files > 0:
            st.info(f"🆕 {actual_new_files} new or updated files detected!")
            
            if processing_user and processing_user != current_user:
                st.error(f"❌ Cannot process documents - {processing_user} is currently processing documents")
            else:
                if st.button("🔄 Load New Documents", type="primary"):
                    lock_acquired, lock_user = st.session_state.rag_system._acquire_lock(current_user)
                    if not lock_acquired:
                        st.error(f"❌ Cannot acquire lock - {lock_user} is processing documents")
                    else:
                        try:
                            with st.spinner(f"Loading new documents... ({current_user} is processing)"):
                                st.session_state.rag_system.process_documents(documents_dir)
                                if st.session_state.rag_system.processed:
                                    st.success(f"✅ Successfully processed {len({d.get('file_path') for d in st.session_state.rag_system.documents})} files!")
                                    st.rerun()
                        finally:
                            st.session_state.rag_system._release_lock()
        else:
            st.button("🔄 Load Documents", disabled=True, help="No new documents to load")
            st.caption("All documents are up to date")
        
        st.markdown("---")
        if st.session_state.rag_system.processed:
            st.success("✅ Knowledge base ready!")
        else:
            st.warning("⚠️ Process documents first")
        
        # Statistics section
        st.markdown("### 📊 Quick Stats")
        doc_files = valid_files_list
        total_available = len(valid_files_list)
        
        if total_available > 0:
            from collections import Counter
            extensions = [Path(f).suffix.lower() for f in doc_files]
            ext_counts = Counter(extensions)
            st.caption(f"File types: {dict(ext_counts)}")
        
        if st.session_state.rag_system.processed:
            processed_files = len({d.get('file_path') for d in st.session_state.rag_system.documents if d.get('file_path')})
            st.metric("Total Files Available", total_available)
            st.metric("Files Processed", processed_files)
            if total_available > processed_files:
                st.metric("New Files", total_available - processed_files)
            else:
                st.metric("New Files", 0)
        else:
            st.metric("Total Files Available", total_available)
            st.metric("Files Processed", 0)
            if total_available > 0:
                st.metric("New Files", total_available)

        # Advanced controls (similar to SustainaCube)
        st.markdown("---")
        st.markdown("### ⚙️ Actions")
        colA, colB = st.columns(2)
        with colA:
            if st.button("♻️ Rebuild Index", help="Reprocess all documents and rebuild cache"):
                acquired, who = st.session_state.rag_system._acquire_lock(current_user)
                if not acquired:
                    st.error(f"Cannot rebuild - {who} is processing")
                else:
                    try:
                        st.session_state.rag_system._clear_cache_and_chunks()
                        with st.spinner("Rebuilding index from documents..."):
                            st.session_state.rag_system.process_documents(documents_dir)
                        st.success("Rebuild complete")
                        st.rerun()
                    finally:
                        st.session_state.rag_system._release_lock()
        with colB:
            if st.button("🧹 Clear Cache", help="Remove processed cache and chunks (keeps files)"):
                st.session_state.rag_system._clear_cache_and_chunks()
                st.success("Cache cleared. Click 'Load New Documents' to reprocess.")

        if st.button("🔓 Release Lock", help="Force release processing lock if stuck"):
            st.session_state.rag_system._release_lock()
            st.success("Lock released")

        # Progress and logs
        with st.expander("⏱️ Processing Status", expanded=False):
            prog = st.session_state.rag_system._read_progress()
            if prog:
                st.write(prog)
            else:
                st.caption("No active processing status.")

        with st.expander("🪵 Recent Logs", expanded=False):
            lines = st.session_state.rag_system._get_recent_log_lines(80)
            if lines:
                st.code("".join(lines))
                st.download_button("Download Recent Logs", data="".join(lines), file_name="recent_logs.txt")
            else:
                st.caption("No logs yet.")
        
        # Full log download
        if st.session_state.rag_system.log_file.exists():
            with open(st.session_state.rag_system.log_file, "r", encoding="utf-8") as f:
                full_log_content = f.read()
            st.download_button(
                "📥 Download Full Processing Log", 
                data=full_log_content, 
                file_name="full_processing_log.txt",
                mime="text/plain",
                help="Download complete processing log to analyze skipped files"
            )
        
        # Failed files log download
        failed_log_path = Path(__file__).parent / "failed_files_log.txt"
        if failed_log_path.exists():
            with open(failed_log_path, "r", encoding="utf-8") as f:
                failed_log_content = f.read()
            st.download_button(
                "📥 Download Failed Files Log", 
                data=failed_log_content, 
                file_name="failed_files_log.txt",
                mime="text/plain",
                help="Download detailed log of files that failed to process"
            )
        
        # Skipped files log download (extract from main log)
        skipped_files = st.session_state.rag_system._get_skipped_files()
        if skipped_files:
            skipped_log_content = "SKIPPED FILES LOG\n" + "="*50 + "\n\n"
            for filename in skipped_files:
                skipped_log_content += f"❌ SKIPPED: {filename}\n"
                skipped_log_content += f"   Reason: No text content extracted\n"
                skipped_log_content += f"   {'-'*40}\n"
            
            st.download_button(
                "📥 Download Skipped Files Log", 
                data=skipped_log_content, 
                file_name="skipped_files_log.txt",
                mime="text/plain",
                help="Download list of files that were skipped (no text content)"
            )
        
        # Failed and skipped files section
        failed_files = st.session_state.rag_system._get_failed_files()
        skipped_files = st.session_state.rag_system._get_skipped_files()
        
        # Debug: Show what we found
        with st.expander("🔍 Debug Info", expanded=False):
            st.write(f"Failed files found: {len(failed_files)}")
            st.write(f"Skipped files found: {len(skipped_files)}")
            if skipped_files:
                st.write("Skipped files:", skipped_files)
        
        if failed_files or skipped_files:
            with st.expander("⚠️ Excluded Files", expanded=False):
                if failed_files:
                    st.warning(f"❌ {len(failed_files)} files failed to process")
                    for failed_file in failed_files:
                        st.text(f"• {failed_file}")
                
                if skipped_files:
                    st.info(f"⏭️ {len(skipped_files)} files skipped (duplicates/unchanged)")
                    for skipped_file in skipped_files:
                        st.text(f"• {skipped_file}")
                
                # Download failed files log
                failed_log_path = Path(__file__).parent / "failed_files_log.txt"
                if failed_log_path.exists():
                    with open(failed_log_path, "r", encoding="utf-8") as f:
                        failed_log_content = f.read()
                    st.download_button(
                        "Download Failed Files Log", 
                        data=failed_log_content, 
                        file_name="failed_files_log.txt",
                        mime="text/plain"
                    )
        else:
            st.success("✅ No excluded files")

    # Main interface
    col1, col2 = st.columns([3, 1])

    with col1:
        st.header("💬 Ask a Question")

        # Manage question state
        if 'question_input' not in st.session_state:
            st.session_state.question_input = ""
        if 'auto_run' not in st.session_state:
            st.session_state.auto_run = False
        # Apply any pending sample selection before rendering the widget
        if st.session_state.get('apply_pending_question', False):
            st.session_state.question_input = st.session_state.get('pending_question', '')
            st.session_state.auto_run = True
            st.session_state.apply_pending_question = False

        # Question input
        question = st.text_area(
            "Enter your PU question:",
            placeholder="e.g., What are the main factors affecting foam density in flexible PU foams?",
            height=100,
            key="question_input",
        )
        # Bound to session via key="question_input"

        def run_query(q: str, mode: str):
            with st.spinner("Searching knowledge base and generating answer..."):
                original_assistant_id = st.session_state.rag_system.assistant_id

                # Knowledge Base only
                if mode == "Knowledge Base only":
                    st.session_state.rag_system.assistant_id = None
                    search_results = st.session_state.rag_system.search_documents(q)
                    if not search_results:
                        answer, sources = ("No relevant information found in the knowledge base.", [])
                    else:
                        answer = st.session_state.rag_system.generate_answer(q, search_results)
                        sources = [
                            {
                                'filename': r['filename'],
                                'similarity_score': r['similarity'],
                                'matched_words': r['matched_words']
                            } for r in search_results
                        ]
                    st.session_state.rag_system.assistant_id = original_assistant_id

                # Assistant only
                elif mode == "Assistant only":
                    if not original_assistant_id:
                        answer, sources = ("Assistant ID not configured. Please set OPENAI_ASSISTANT_ID.", [])
                    else:
                        answer, sources = st.session_state.rag_system.answer_question(q)

                # Hybrid: KB then Assistant refine
                else:
                    # Step 1: Local KB retrieval
                    st.session_state.rag_system.assistant_id = None
                    search_results = st.session_state.rag_system.search_documents(q)
                    kb_answer = ""
                    if search_results:
                        kb_answer = st.session_state.rag_system.generate_answer(q, search_results)
                    st.session_state.rag_system.assistant_id = original_assistant_id

                    # If no KB context found, fall back to Assistant-only answer
                    if not search_results and original_assistant_id:
                        answer, sources = st.session_state.rag_system.answer_question(q)
                    else:
                        # Step 2: Assistant refinement
                        if not original_assistant_id:
                            answer = kb_answer or "Assistant ID not configured and no KB answer available."
                            sources = []
                        else:
                            try:
                                # Build refinement prompt
                                context_text = "\n\n".join([r['text'] for r in (search_results or [])][:5])
                                refine_prompt = (
                                    "You are a polyurethane expert. Refine the following draft answer. "
                                    "If the draft lacks substance, improve it using your tools and knowledge. "
                                    "Preserve a professional, structured format.\n\n"
                                    f"Question: {q}\n\n"
                                    f"Context (from local KB):\n{context_text}\n\n"
                                    f"Draft answer:\n{kb_answer}"
                                )
                                client = st.session_state.rag_system.openai_client
                                thread = client.beta.threads.create()
                                client.beta.threads.messages.create(
                                    thread_id=thread.id,
                                    role="user",
                                    content=refine_prompt
                                )
                                run = client.beta.threads.runs.create(
                                    thread_id=thread.id,
                                    assistant_id=original_assistant_id
                                )
                                start = time.time()
                                while True:
                                    r = client.beta.threads.runs.retrieve(thread_id=thread.id, run_id=run.id)
                                    if r.status in ["completed", "failed", "cancelled", "expired"]:
                                        break
                                    time.sleep(0.4)
                                    if time.time() - start > 120:
                                        break
                                if r.status == "completed":
                                    msgs = client.beta.threads.messages.list(thread_id=thread.id)
                                    parts = []
                                    for m in msgs.data:
                                        if m.role == 'assistant':
                                            for c in m.content:
                                                if getattr(c, 'type', '') == 'text':
                                                    parts.append(c.text.value)
                                            break
                                    answer = "\n\n".join(parts) if parts else (kb_answer or "")
                                else:
                                    answer = kb_answer or f"Assistant run status: {r.status}"
                                # Collect sources (basic)
                                sources = []
                                if search_results:
                                    for result in search_results[:5]:
                                        sources.append({'filename': result['filename'], 'similarity_score': result['similarity']})
                            except Exception as e:
                                answer = kb_answer or f"Hybrid refinement error: {e}"
                                sources = []

            st.markdown("### 📋 Answer")
            st.caption(f"Answer source: {mode}")
            st.markdown(answer)

            # Convert markdown to HTML and provide professional styling
            import re
            import time

            html_answer = answer
            html_answer = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^#### (.+)$', r'<h4>\1</h4>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'【[^】]+】', '', html_answer)
            html_answer = re.sub(r'<h3>Source References</h3>.*?(?=<h3>|$)', '', html_answer, flags=re.DOTALL)
            html_answer = re.sub(r'^- (.+)$', r'<li>\1</li>', html_answer, flags=re.MULTILINE)
            html_answer = re.sub(r'(<li>.*</li>)(?:\s*<li>.*</li>)*', lambda m: f'<ul>{m.group(0)}</ul>', html_answer, flags=re.DOTALL)
            html_answer = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html_answer)
            html_answer = html_answer.replace('\n', '<br>')

            source_refs = []
            if sources:
                for source in sources:
                    src_name = source.get('filename') if isinstance(source, dict) else str(source)
                    score = source.get('similarity_score') if isinstance(source, dict) else None
                    if isinstance(score, (int, float)):
                        source_refs.append(f'<li>{src_name} (Relevance: {score:.3f})</li>')
                    else:
                        source_refs.append(f'<li>{src_name}</li>')
            source_list = f'<ul>{"".join(source_refs)}</ul>' if source_refs else '<p><em>No specific sources referenced.</em></p>'

            html_content = f"""
<!doctype html>
<html>
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <style>
    body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif; line-height: 1.7; color: #2c3e50; margin: 0; padding: 40px; background: #f8f9fa; }}
    .container {{ max-width: 800px; margin: 0 auto; background: white; padding: 40px; border-radius: 8px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
    h1 {{ color: #2c3e50; font-size: 28px; font-weight: 700; margin: 0 0 20px 0; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
    h2 {{ color: #34495e; font-size: 22px; font-weight: 600; margin: 30px 0 15px 0; border-left: 4px solid #3498db; padding-left: 15px; }}
    h3 {{ color: #34495e; font-size: 18px; font-weight: 600; margin: 25px 0 12px 0; }}
    h4 {{ color: #34495e; font-size: 16px; font-weight: 600; margin: 20px 0 10px 0; }}
    ul {{ margin: 15px 0; padding-left: 20px; }}
    li {{ margin: 8px 0; line-height: 1.6; }}
    p {{ margin: 15px 0; line-height: 1.7; }}
    strong {{ color: #2c3e50; font-weight: 600; }}
    .header {{ text-align: center; margin-bottom: 30px; padding-bottom: 20px; border-bottom: 2px solid #ecf0f1; }}
    .header h1 {{ border: none; margin: 0; color: #2c3e50; }}
    .timestamp {{ color: #7f8c8d; font-size: 14px; margin-top: 10px; }}
    .question-section {{ background: #f8f9fa; padding: 20px; border-radius: 6px; margin: 20px 0; border-left: 4px solid #3498db; }}
    .question-text {{ font-size: 16px; color: #2c3e50; font-weight: 500; margin: 10px 0 0 0; line-height: 1.6; }}
  </style>
  <title>PU ExpertCenter Response</title>
</head>
<body>
  <div class="container">
    <div class="header"><h1>🧪 PU ExpertCenter Response</h1><div class="timestamp">Generated on {time.strftime('%B %d, %Y at %I:%M %p')}</div></div>
    <div class="question-section"><h2>❓ Question</h2><p class="question-text">{q}</p></div>
    <div class="content">{html_answer}</div>
  </div>
</body>
</html>
"""
            st.download_button(
                label="Copy Answer (HTML)",
                data=html_content,
                file_name="pu_expertcenter_answer.html",
                mime="text/html",
                key="download_html"
            )
            st.download_button(
                label="Copy Answer (Text)",
                data=answer,
                file_name="pu_expertcenter_answer.txt",
                mime="text/plain",
                key="download_text"
            )

        # Assistant toggle
        st.markdown("---")
        st.markdown("### 🤖 Answer Mode")
        mode = st.radio(
            "Choose how to generate the answer",
            options=["Knowledge Base only", "Assistant only", "Hybrid (KB → Assistant refine)"],
            index=0,
            help="Hybrid first builds a KB answer, then asks the Assistant to refine it."
        )
        # Normalize label for logic
        if mode.startswith("Hybrid"):
            selected_mode = "Hybrid"
        else:
            selected_mode = mode

        if st.button("🔍 Get Answer", type="primary"):
            if question.strip():
                run_query(question, selected_mode)
            else:
                st.warning("Please enter a question.")

        if st.session_state.auto_run and st.session_state.question_input.strip():
            run_query(st.session_state.question_input, selected_mode)
            st.session_state.auto_run = False

    with col2:
        st.markdown("### 💡 Sample Questions")
        sample_questions = [
            # General domain
            "What are the main types of polyurethane catalysts?",
            "What factors affect the density of rigid PU foam?",
            "What are typical compression set values for automotive foams?",
            "How does temperature affect gel time in PU systems?",
            # KB-targeted (should hit specific docs)
            "According to the Troubleshooting guide - Laader Berg, what causes foam collapse?",
            "From the Troubleshooting guide - Laader Berg, list common root causes of shrinkage.",
            "In Dow Polyurethanes Flexible Foams, how does water level impact density and cell structure?",
            "What key market trends does EUROPUR Market Report FY 2024 highlight?",
            "From Safety-Guidelines-2023-1, what PPE is required for handling TDI/MDI?",
            "What are the main findings in 'Handbook-of-plastic-foams' about compression behavior?",
            "Summarize key ESG metrics mentioned in hennecke_esg_key_mertrics_2024.",
            "What troubleshooting steps address voids/porosity in slabstock foams per Laader Berg?",
        ]

        for q in sample_questions:
            if st.button(q, key=f"sample_btn_{hash(q)}"):
                st.session_state.pending_question = q
                st.session_state.apply_pending_question = True
                st.rerun()

if __name__ == "__main__":
    main()
